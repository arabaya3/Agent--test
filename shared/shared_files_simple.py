import requests
import base64
import json
import webbrowser
import os
import subprocess
import sys
import io
import time
from dotenv import load_dotenv
from .auth import get_access_token

load_dotenv()

allowed_extensions = [
    '.pdf', '.docx', '.doc', '.txt', '.xlsx', '.xls', '.ppt', '.pptx',
    '.zip', '.rar', '.7z', '.csv', '.rtf', '.odt', '.ods', '.odp',
    '.html', '.htm', '.xml', '.json', '.log', '.md', '.tex'
]

text_extensions = [
    '.txt', '.csv', '.log', '.md', '.tex', '.html', '.htm', '.xml', '.json'
]

def is_allowed_file(filename):
    file_ext = os.path.splitext(filename.lower())[1]
    return file_ext in allowed_extensions

def is_text_file(filename):
    file_ext = os.path.splitext(filename.lower())[1]
    return file_ext in text_extensions

def get_conversation_messages(conversation_id, headers, user_id=None):
    """Get all messages in a conversation thread"""
    base_url = "https://graph.microsoft.com/v1.0"
    
    # For application permissions, we need to specify a user ID
    if user_id:
        user_prefix = f"users/{user_id}"
    else:
        user_prefix = "me"
    
    # 1. Try $search (if enabled)
    search_url = f"{base_url}/{user_prefix}/messages?$search=\"conversationId:{conversation_id}\""
    print(f"[DEBUG] Requesting ($search): {search_url}")
    try:
        response = requests.get(search_url, headers=headers, timeout=30)
        if response.status_code == 200:
            messages = response.json().get("value", [])
            if messages:
                print(f"[DEBUG] Found {len(messages)} messages using $search.")
                return messages
            else:
                print(f"[DEBUG] No messages found using $search for conversationId: {conversation_id}")
        else:
            print(f"[DEBUG] Status {response.status_code}: {response.text}")
    except requests.exceptions.RequestException as e:
        print(f"Error retrieving conversation ($search) {conversation_id}: {e}")
    
    # 2. Try msgfolderroot (AllItems) with smaller limit
    allitems_url = (
        f"{base_url}/{user_prefix}/mailFolders/msgfolderroot/messages?"
        f"$filter=conversationId eq '{conversation_id}'&$orderby=sentDateTime asc&$top=50"
    )
    print(f"[DEBUG] Requesting (msgfolderroot): {allitems_url}")
    try:
        response = requests.get(allitems_url, headers=headers, timeout=30)
        if response.status_code == 200:
            messages = response.json().get("value", [])
            if messages:
                print(f"[DEBUG] Found {len(messages)} messages in msgfolderroot.")
                return messages
            else:
                print(f"[DEBUG] No messages found in msgfolderroot for conversationId: {conversation_id}")
        else:
            print(f"[DEBUG] Status {response.status_code}: {response.text}")
    except requests.exceptions.RequestException as e:
        print(f"Error retrieving conversation (msgfolderroot) {conversation_id}: {e}")
    
    # 3. Try inbox with smaller limit
    inbox_url = (
        f"{base_url}/{user_prefix}/mailFolders/inbox/messages?"
        f"$filter=conversationId eq '{conversation_id}'&$orderby=sentDateTime asc&$top=50"
    )
    print(f"[DEBUG] Requesting (inbox): {inbox_url}")
    try:
        response = requests.get(inbox_url, headers=headers, timeout=30)
        if response.status_code == 200:
            messages = response.json().get("value", [])
            if messages:
                print(f"[DEBUG] Found {len(messages)} messages in inbox.")
                return messages
            else:
                print(f"[DEBUG] No messages found in inbox for conversationId: {conversation_id}")
        else:
            print(f"[DEBUG] Status {response.status_code}: {response.text}")
    except requests.exceptions.RequestException as e:
        print(f"Error retrieving conversation (inbox) {conversation_id}: {e}")
    
    # 4. Limited fallback: Fetch recent messages only (max 1000)
    print(f"[DEBUG] Trying limited fallback: fetching recent messages only...")
    all_msgs = []
    next_url = f"{base_url}/{user_prefix}/messages?$top=100&$orderby=receivedDateTime desc"
    message_count = 0
    max_messages = 1000  # Limit to prevent excessive API calls
    
    while next_url and message_count < max_messages:
        print(f"[DEBUG] Requesting (limited): {next_url}")
        try:
            response = requests.get(next_url, headers=headers, timeout=30)
            if response.status_code == 200:
                data = response.json()
                batch = data.get("value", [])
                filtered = [m for m in batch if m.get("conversationId") == conversation_id]
                all_msgs.extend(filtered)
                message_count += len(batch)
                next_url = data.get("@odata.nextLink")
                if next_url:
                    time.sleep(0.1)  # Reduced delay
            else:
                print(f"[DEBUG] Status {response.status_code}: {response.text}")
                break
        except requests.exceptions.RequestException as e:
            print(f"Error retrieving messages for client-side filtering: {e}")
            break
    
    if all_msgs:
        print(f"[DEBUG] Found {len(all_msgs)} messages by limited fallback (searched {message_count} recent messages).")
        all_msgs.sort(key=lambda m: m.get("sentDateTime", ""))
        return all_msgs
    
    print(f"[DEBUG] All attempts failed for conversationId: {conversation_id}")
    print(f"[DEBUG] Conversation not found in recent {message_count} messages.")
    return []

def extract_pdf_text_from_bytes(content_bytes):
    try:
        import PyPDF2
        pdf_file = io.BytesIO(content_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except ImportError:
        print("PyPDF2 not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "PyPDF2"])
            import PyPDF2
            pdf_file = io.BytesIO(content_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        return None

def extract_docx_text_from_bytes(content_bytes):
    try:
        from docx import Document
        doc_file = io.BytesIO(content_bytes)
        doc = Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except ImportError:
        print("python-docx not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx"])
            from docx import Document
            doc_file = io.BytesIO(content_bytes)
            doc = Document(doc_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting DOCX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return None

def extract_xlsx_text_from_bytes(content_bytes):
    try:
        import pandas as pd
        excel_file = io.BytesIO(content_bytes)
        df = pd.read_excel(excel_file)
        text = df.to_string(index=False)
        return text
    except ImportError:
        print("pandas not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "pandas", "openpyxl"])
            import pandas as pd
            excel_file = io.BytesIO(content_bytes)
            df = pd.read_excel(excel_file)
            text = df.to_string(index=False)
            return text
        except Exception as e:
            print(f"Error extracting XLSX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting XLSX text: {e}")
        return None

def extract_pptx_text_from_bytes(content_bytes):
    try:
        from pptx import Presentation
        pptx_file = io.BytesIO(content_bytes)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except ImportError:
        print("python-pptx not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            from pptx import Presentation
            pptx_file = io.BytesIO(content_bytes)
            prs = Presentation(pptx_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PPTX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return None

def get_access_token():
    """Get access token using application permissions (client credentials flow)"""
    from .auth import get_access_token as get_app_token
    return get_app_token()

def get_email_by_id(email_id, headers, user_id=None):
    if user_id:
        email_url = f'https://graph.microsoft.com/v1.0/users/{user_id}/messages/{email_id}'
    else:
        email_url = f'https://graph.microsoft.com/v1.0/me/messages/{email_id}'
    response = requests.get(email_url, headers=headers)
    
    if response.status_code == 200:
        return response.json()
    else:
        print(f"Error getting email: {response.status_code}")
        print(f"Error: {response.text}")
        return None

def get_attachments_for_email(email_id, headers, user_id=None):
    if user_id:
        att_url = f'https://graph.microsoft.com/v1.0/users/{user_id}/messages/{email_id}/attachments'
    else:
        att_url = f'https://graph.microsoft.com/v1.0/me/messages/{email_id}/attachments'
    response = requests.get(att_url, headers=headers)
    
    if response.status_code == 200:
        return response.json().get('value', [])
    else:
        print(f"Error getting attachments: {response.status_code}")
        return []

def print_file_content_from_bytes(attachment):
    filename = attachment['name']
    content_bytes = base64.b64decode(attachment['contentBytes'])
    file_ext = os.path.splitext(filename.lower())[1]
    
    if is_text_file(filename):
        try:
            content = content_bytes.decode('utf-8')
            if len(content.strip()) == 0:
                print(f"\nFile: {filename}")
                print("="*60)
                print("File is empty or contains no readable text.")
                print("="*60)
                return
            
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        except UnicodeDecodeError:
            try:
                content = content_bytes.decode('latin-1')
                if len(content.strip()) == 0:
                    print(f"\nFile: {filename}")
                    print("="*60)
                    print("File is empty or contains no readable text.")
                    print("="*60)
                    return
                
                print(f"\nContent of {filename}:")
                print("="*60)
                print(content)
                print("="*60)
            except:
                print(f"\nFile: {filename}")
                print("="*60)
                print("Cannot read this file as text.")
                print("="*60)
        except Exception as e:
            print(f"\nFile: {filename}")
            print("="*60)
            print(f"Error reading file: {e}")
            print("="*60)
    
    elif file_ext == '.pdf':
        print(f"\nExtracting text from PDF: {filename}")
        content = extract_pdf_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from PDF.")
            print("="*60)
    
    elif file_ext == '.docx':
        print(f"\nExtracting text from DOCX: {filename}")
        content = extract_docx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from DOCX.")
            print("="*60)
    
    elif file_ext in ['.xlsx', '.xls']:
        print(f"\nExtracting text from Excel: {filename}")
        content = extract_xlsx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from Excel file.")
            print("="*60)
    
    elif file_ext in ['.pptx', '.ppt']:
        print(f"\nExtracting text from PowerPoint: {filename}")
        content = extract_pptx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from PowerPoint file.")
            print("="*60)
    
    else:
        print(f"\nFile: {filename}")
        print("="*60)
        print("This file type is not supported for text extraction.")
        print("="*60)

def get_all_attachments_from_conversation(email_id, headers):
    """Get all attachments from the original email and all replies in the conversation"""
    print("============================================================")
    print("Email Attachment Viewer (with Conversation Thread)")
    print("============================================================")
    
    # Get the original email
    email_data = get_email_by_id(email_id, headers)
    if not email_data:
        print("Email not found!")
        return []
    
    print(f"\nOriginal Email Details:")
    print(f"Subject: {email_data.get('subject', 'No subject')}")
    print(f"From: {email_data.get('from', {}).get('emailAddress', {}).get('address', 'Unknown')}")
    print(f"Date: {email_data.get('receivedDateTime', 'Unknown')}")
    
    # Get conversation ID
    conversation_id = email_data.get("conversationId")
    if not conversation_id:
        print("No conversation ID found for this email.")
        return []
    
    # Get all messages in the conversation
    print(f"\nRetrieving conversation thread...")
    conversation_messages = get_conversation_messages(conversation_id, headers)
    
    if not conversation_messages:
        print("No conversation messages found.")
        return []
    
    # Combine original email and conversation messages, remove duplicates
    orig_id = str(email_data.get("id")).lower()
    all_messages = [email_data]  # Start with original email
    
    for msg in conversation_messages:
        if str(msg.get("id")).lower() != orig_id:
            all_messages.append(msg)
    
    # Sort by received date
    all_messages.sort(key=lambda m: m.get("receivedDateTime", ""))
    
    print(f"\nFound {len(all_messages)} messages in conversation thread")
    
    # Collect all attachments from all messages
    all_attachments = []
    attachment_sources = []  # Track which email each attachment comes from
    
    for i, message in enumerate(all_messages, 1):
        message_id = message.get("id")
        message_subject = message.get("subject", "No subject")
        message_from = message.get("from", {}).get("emailAddress", {}).get("address", "Unknown")
        message_date = message.get("receivedDateTime", "Unknown")
        
        print(f"\nProcessing message {i}/{len(all_messages)}: {message_subject}")
        print(f"From: {message_from}")
        print(f"Date: {message_date}")
        
        attachments = get_attachments_for_email(message_id, headers)
        
        if attachments:
            print(f"Found {len(attachments)} attachments in this message")
            
            for attachment in attachments:
                if attachment['@odata.type'] == '#microsoft.graph.fileAttachment':
                    if is_allowed_file(attachment['name']):
                        all_attachments.append(attachment)
                        attachment_sources.append({
                            'email_subject': message_subject,
                            'email_from': message_from,
                            'email_date': message_date,
                            'message_number': i
                        })
                        print(f"  ✓ Added: {attachment['name']}")
                    else:
                        print(f"  - Skipped: {attachment['name']} (not a document)")
        else:
            print("No attachments in this message")
    
    print(f"\n============================================================")
    print(f"Total attachments found in conversation: {len(all_attachments)}")
    print("============================================================")
    
    return all_attachments, attachment_sources

def main():
    if not client_id or not tenant_id:
        print("Error: CLIENT_ID and TENANT_ID must be set in .env file")
        print("Please create a .env file with the following content:")
        print("CLIENT_ID=your_client_id_here")
        print("TENANT_ID=your_tenant_id_here")
        return
    
    access_token = get_access_token()
    
    if not access_token:
        print("Failed to get access token")
        return
    
    headers = {
        'Authorization': f'Bearer {access_token}',
        'Content-Type': 'application/json'
    }
    
    print("\n" + "="*60)
    print("Email Attachment Viewer (Conversation Thread)")
    print("="*60)
    
    email_id = input("Enter the email ID: ").strip()
    
    if not email_id:
        print("Email ID is required!")
        return
    
    print(f"\nFetching email with ID: {email_id}")
    
    # Get all attachments from conversation thread
    all_attachments, attachment_sources = get_all_attachments_from_conversation(email_id, headers)
    
    if not all_attachments:
        print("\nNo document attachments found in the conversation thread.")
        return
    
    print(f"\nDocument attachments available: {len(all_attachments)}")
    print("Note: The script will extract text directly from attachments without saving files.")
    print("Required libraries will be installed automatically if needed.")
    
    # Display all attachments with their source information
    print(f"\nAll attachments in conversation:")
    print("="*80)
    for i, (attachment, source) in enumerate(zip(all_attachments, attachment_sources), 1):
        file_type = "Text" if is_text_file(attachment['name']) else "Document"
        print(f"{i:2d}. {attachment['name']} ({file_type})")
        print(f"     From: {source['email_from']}")
        print(f"     Subject: {source['email_subject']}")
        print(f"     Date: {source['email_date']}")
        print()
    
    while True:
        try:
            choice = input(f"\nSelect a file to view (1-{len(all_attachments)}) or 'q' to quit: ").strip()
            
            if choice.lower() == 'q':
                print("Goodbye!")
                break
            
            file_index = int(choice) - 1
            
            if 0 <= file_index < len(all_attachments):
                selected_attachment = all_attachments[file_index]
                source_info = attachment_sources[file_index]
                
                print(f"\nViewing attachment from:")
                print(f"Email: {source_info['email_subject']}")
                print(f"From: {source_info['email_from']}")
                print(f"Date: {source_info['email_date']}")
                
                print_file_content_from_bytes(selected_attachment)
                
                view_another = input("\nView another file? (y/n): ").strip().lower()
                if view_another != 'y':
                    break
            else:
                print(f"Please enter a number between 1 and {len(all_attachments)}")
        except ValueError:
            print("Please enter a valid number")
        except KeyboardInterrupt:
            print("\nGoodbye!")
            break

if __name__ == "__main__":
    main() 