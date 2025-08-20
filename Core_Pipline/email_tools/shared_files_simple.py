import requests
import base64
import json
import webbrowser
import os
import subprocess
import sys
import io
import time
from dotenv import load_dotenv
try:
    from shared.auth import get_access_token
except ImportError:
    import sys, os
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if base_dir not in sys.path:
        sys.path.append(base_dir)
    from shared.auth import get_access_token

load_dotenv()

allowed_extensions = [
    '.pdf', '.docx', '.doc', '.txt', '.xlsx', '.xls', '.ppt', '.pptx',
    '.zip', '.rar', '.7z', '.csv', '.rtf', '.odt', '.ods', '.odp',
    '.html', '.htm', '.xml', '.json', '.log', '.md', '.tex'
]

text_extensions = [
    '.txt', '.csv', '.log', '.md', '.tex', '.html', '.htm', '.xml', '.json'
]

def is_allowed_file(filename):
    file_ext = os.path.splitext(filename.lower())[1]
    return file_ext in allowed_extensions

def is_text_file(filename):
    file_ext = os.path.splitext(filename.lower())[1]
    return file_ext in text_extensions

def get_conversation_messages(conversation_id, headers, user_id=None):
    base_url = "https://graph.microsoft.com/v1.0"
    target_user = user_id or os.getenv("DEFAULT_USER_ID")
    if not target_user:
        print("Error: No user ID provided and DEFAULT_USER_ID not set")
        return []
    user_prefix = f"users/{target_user}"

    search_url = f"{base_url}/{user_prefix}/messages?$search=\"conversationId:{conversation_id}\""
    try:
        response = requests.get(search_url, headers=headers, timeout=30)
        if response.status_code == 200:
            messages = response.json().get("value", [])
            if messages:
                return messages
    except requests.exceptions.RequestException:
        pass

    allitems_url = (
        f"{base_url}/{user_prefix}/mailFolders/msgfolderroot/messages?"
        f"$filter=conversationId eq '{conversation_id}'&$orderby=sentDateTime asc&$top=50"
    )
    try:
        response = requests.get(allitems_url, headers=headers, timeout=30)
        if response.status_code == 200:
            messages = response.json().get("value", [])
            if messages:
                return messages
    except requests.exceptions.RequestException:
        pass

    inbox_url = (
        f"{base_url}/{user_prefix}/mailFolders/inbox/messages?"
        f"$filter=conversationId eq '{conversation_id}'&$orderby=sentDateTime asc&$top=50"
    )
    try:
        response = requests.get(inbox_url, headers=headers, timeout=30)
        if response.status_code == 200:
            messages = response.json().get("value", [])
            if messages:
                return messages
    except requests.exceptions.RequestException:
        pass

    all_msgs = []
    next_url = f"{base_url}/{user_prefix}/messages?$top=100&$orderby=receivedDateTime desc"
    message_count = 0
    max_messages = 1000
    while next_url and message_count < max_messages:
        try:
            response = requests.get(next_url, headers=headers, timeout=30)
            if response.status_code == 200:
                data = response.json()
                batch = data.get("value", [])
                filtered = [m for m in batch if m.get("conversationId") == conversation_id]
                all_msgs.extend(filtered)
                message_count += len(batch)
                next_url = data.get("@odata.nextLink")
                if next_url:
                    time.sleep(0.1)
            else:
                break
        except requests.exceptions.RequestException:
            break

    if all_msgs:
        all_msgs.sort(key=lambda m: m.get("sentDateTime", ""))
        return all_msgs

    return []

def extract_pdf_text_from_bytes(content_bytes):
    try:
        import PyPDF2
        pdf_file = io.BytesIO(content_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except ImportError:
        print("PyPDF2 not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "PyPDF2"])
            import PyPDF2
            pdf_file = io.BytesIO(content_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        return None

def extract_docx_text_from_bytes(content_bytes):
    try:
        from docx import Document
        doc_file = io.BytesIO(content_bytes)
        doc = Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except ImportError:
        print("python-docx not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx"])
            from docx import Document
            doc_file = io.BytesIO(content_bytes)
            doc = Document(doc_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting DOCX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return None

def extract_xlsx_text_from_bytes(content_bytes):
    try:
        import pandas as pd
        excel_file = io.BytesIO(content_bytes)
        df = pd.read_excel(excel_file)
        text = df.to_string(index=False)
        return text
    except ImportError:
        print("pandas not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "pandas", "openpyxl"])
            import pandas as pd
            excel_file = io.BytesIO(content_bytes)
            df = pd.read_excel(excel_file)
            text = df.to_string(index=False)
            return text
        except Exception as e:
            print(f"Error extracting XLSX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting XLSX text: {e}")
        return None

def extract_pptx_text_from_bytes(content_bytes):
    try:
        from pptx import Presentation
        pptx_file = io.BytesIO(content_bytes)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except ImportError:
        print("python-pptx not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            from pptx import Presentation
            pptx_file = io.BytesIO(content_bytes)
            prs = Presentation(pptx_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PPTX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return None

def get_email_by_id(email_id, headers, user_id=None):
    base_url = 'https://graph.microsoft.com/v1.0'
    target_user = user_id or os.getenv("DEFAULT_USER_ID")
    if not target_user:
        print("Error: No user ID provided and DEFAULT_USER_ID not set")
        return None
    user_prefix = f'users/{target_user}'
    email_url = f'{base_url}/{user_prefix}/messages/{email_id}'
    response = requests.get(email_url, headers=headers)
    if response.status_code == 200:
        return response.json()
    else:
        print(f"Error getting email: {response.status_code}")
        print(f"Error: {response.text}")
        return None

def get_attachments_for_email(email_id, headers, user_id=None):
    base_url = 'https://graph.microsoft.com/v1.0'
    target_user = user_id or os.getenv("DEFAULT_USER_ID")
    if not target_user:
        print("Error: No user ID provided and DEFAULT_USER_ID not set")
        return []
    user_prefix = f'users/{target_user}'
    att_url = f'{base_url}/{user_prefix}/messages/{email_id}/attachments'
    response = requests.get(att_url, headers=headers)
    if response.status_code == 200:
        return response.json().get('value', [])
    else:
        print(f"Error getting attachments: {response.status_code}")
        return []

def print_file_content_from_bytes(attachment):
    filename = attachment['name']
    content_bytes = base64.b64decode(attachment['contentBytes'])
    file_ext = os.path.splitext(filename.lower())[1]

    if is_text_file(filename):
        try:
            content = content_bytes.decode('utf-8')
            if len(content.strip()) == 0:
                print(f"\nFile: {filename}")
                print("="*60)
                print("File is empty or contains no readable text.")
                print("="*60)
                return
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        except UnicodeDecodeError:
            try:
                content = content_bytes.decode('latin-1')
                if len(content.strip()) == 0:
                    print(f"\nFile: {filename}")
                    print("="*60)
                    print("File is empty or contains no readable text.")
                    print("="*60)
                    return
                print(f"\nContent of {filename}:")
                print("="*60)
                print(content)
                print("="*60)
            except Exception:
                print(f"\nFile: {filename}")
                print("="*60)
                print("Cannot read this file as text.")
                print("="*60)
        except Exception as e:
            print(f"\nFile: {filename}")
            print("="*60)
            print(f"Error reading file: {e}")
            print("="*60)
    elif file_ext == '.pdf':
        print(f"\nExtracting text from PDF: {filename}")
        content = extract_pdf_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from PDF.")
            print("="*60)
    elif file_ext == '.docx':
        print(f"\nExtracting text from DOCX: {filename}")
        content = extract_docx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from DOCX.")
            print("="*60)
    elif file_ext in ['.xlsx', '.xls']:
        print(f"\nExtracting text from Excel: {filename}")
        content = extract_xlsx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from Excel file.")
            print("="*60)
    elif file_ext in ['.pptx', '.ppt']:
        print(f"\nExtracting text from PowerPoint: {filename}")
        content = extract_pptx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from PowerPoint file.")
            print("="*60)
    else:
        print(f"\nFile: {filename}")
        print("="*60)
        print("This file type is not supported for text extraction.")
        print("="*60)

def get_all_attachments_from_conversation(email_id, headers, user_id=None):
    print("============================================================")
    print("Email Attachment Viewer (with Conversation Thread)")
    print("============================================================")

    email_data = get_email_by_id(email_id, headers, user_id)
    if not email_data:
        print("Email not found!")
        return [], []

    print(f"\nOriginal Email Details:")
    print(f"Subject: {email_data.get('subject', 'No subject')}")
    print(f"From: {email_data.get('from', {}).get('emailAddress', {}).get('address', 'Unknown')}")
    print(f"Date: {email_data.get('receivedDateTime', 'Unknown')}")

    conversation_id = email_data.get("conversationId")
    if not conversation_id:
        print("No conversation ID found for this email.")
        return [], []

    print(f"\nRetrieving conversation thread...")
    conversation_messages = get_conversation_messages(conversation_id, headers, user_id)
    if not conversation_messages:
        print("No conversation messages found.")
        return [], []

    orig_id = str(email_data.get("id")).lower()
    all_messages = [email_data]
    for msg in conversation_messages:
        if str(msg.get("id")).lower() != orig_id:
            all_messages.append(msg)
    all_messages.sort(key=lambda m: m.get("receivedDateTime", ""))

    all_attachments = []
    attachment_sources = []
    for i, message in enumerate(all_messages, 1):
        message_id = message.get("id")
        message_subject = message.get("subject", "No subject")
        message_from = message.get("from", {}).get("emailAddress", {}).get("address", "Unknown")
        message_date = message.get("receivedDateTime", "Unknown")
        attachments = get_attachments_for_email(message_id, headers, user_id)
        if attachments:
            for attachment in attachments:
                if attachment.get('@odata.type') == '#microsoft.graph.fileAttachment':
                    if is_allowed_file(attachment['name']):
                        all_attachments.append(attachment)
                        attachment_sources.append({
                            'email_subject': message_subject,
                            'email_from': message_from,
                            'email_date': message_date,
                            'message_number': i
                        })
    return all_attachments, attachment_sources

def main():
    access_token = get_access_token()
    if not access_token:
        print("Failed to get access token")
        return
    headers = {
        'Authorization': f'Bearer {access_token}',
        'Content-Type': 'application/json'
    }
    print("\n" + "="*60)
    print("📎 EMAIL ATTACHMENT VIEWER")
    print("="*60)
    email_id = input("Enter the email ID: ").strip()
    if not email_id:
        print("Email ID is required!")
        return
    # Get user ID for application permissions
    target_user = os.getenv("DEFAULT_USER_ID")
    if not target_user:
        print("Error: DEFAULT_USER_ID not set. Cannot fetch emails.")
        return
    
    print(f"\nFetching email with ID: {email_id}")
    all_attachments, attachment_sources = get_all_attachments_from_conversation(email_id, headers, target_user)
    if not all_attachments:
        print("\nNo document attachments found in the conversation thread.")
        return
    print(f"\n📎 ATTACHMENTS FOUND: {len(all_attachments)}")
    print("=" * 60)
    for i, (attachment, source) in enumerate(zip(all_attachments, attachment_sources), 1):
        file_type = "Text" if is_text_file(attachment['name']) else "Document"
        print(f"\n📄 {i}. {attachment['name']} ({file_type})")
        print(f"   👤 From: {source['email_from']}")
        print(f"   📋 Subject: {source['email_subject']}")
        print(f"   📅 Date: {source['email_date']}")
        print("─" * 40)
    while True:
        try:
            choice = input(f"\nSelect a file to view (1-{len(all_attachments)}) or 'q' to quit: ").strip()
            if choice.lower() == 'q':
                print("Goodbye!")
                break
            file_index = int(choice) - 1
            if 0 <= file_index < len(all_attachments):
                selected_attachment = all_attachments[file_index]
                source_info = attachment_sources[file_index]
                print(f"\nViewing attachment from:")
                print(f"Email: {source_info['email_subject']}")
                print(f"From: {source_info['email_from']}")
                print(f"Date: {source_info['email_date']}")
                print_file_content_from_bytes(selected_attachment)
                view_another = input("\nView another file? (y/n): ").strip().lower()
                if view_another != 'y':
                    break
            else:
                print(f"Please enter a number between 1 and {len(all_attachments)}")
        except ValueError:
            print("Please enter a valid number")
        except KeyboardInterrupt:
            print("\nGoodbye!")
            break

if __name__ == "__main__":
    main()


