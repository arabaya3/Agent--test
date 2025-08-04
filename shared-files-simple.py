from msal import PublicClientApplication
import requests
import base64
import json
import webbrowser
import os
import subprocess
import sys
import io
from dotenv import load_dotenv

load_dotenv()

client_id = os.getenv('CLIENT_ID')
tenant_id = os.getenv('TENANT_ID')
authority = f'https://login.microsoftonline.com/{tenant_id}' if tenant_id else None
scopes = ['https://graph.microsoft.com/Mail.Read', 'https://graph.microsoft.com/User.Read']

allowed_extensions = [
    '.pdf', '.docx', '.doc', '.txt', '.xlsx', '.xls', '.ppt', '.pptx',
    '.zip', '.rar', '.7z', '.csv', '.rtf', '.odt', '.ods', '.odp',
    '.html', '.htm', '.xml', '.json', '.log', '.md', '.tex'
]

text_extensions = [
    '.txt', '.csv', '.log', '.md', '.tex', '.html', '.htm', '.xml', '.json'
]

app = PublicClientApplication(
    client_id=client_id,
    authority=authority,
)

def is_allowed_file(filename):
    file_ext = os.path.splitext(filename.lower())[1]
    return file_ext in allowed_extensions

def is_text_file(filename):
    file_ext = os.path.splitext(filename.lower())[1]
    return file_ext in text_extensions

def extract_pdf_text_from_bytes(content_bytes):
    try:
        import PyPDF2
        pdf_file = io.BytesIO(content_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except ImportError:
        print("PyPDF2 not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "PyPDF2"])
            import PyPDF2
            pdf_file = io.BytesIO(content_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        return None

def extract_docx_text_from_bytes(content_bytes):
    try:
        from docx import Document
        doc_file = io.BytesIO(content_bytes)
        doc = Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except ImportError:
        print("python-docx not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx"])
            from docx import Document
            doc_file = io.BytesIO(content_bytes)
            doc = Document(doc_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting DOCX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return None

def extract_xlsx_text_from_bytes(content_bytes):
    try:
        import pandas as pd
        excel_file = io.BytesIO(content_bytes)
        df = pd.read_excel(excel_file)
        text = df.to_string(index=False)
        return text
    except ImportError:
        print("pandas not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "pandas", "openpyxl"])
            import pandas as pd
            excel_file = io.BytesIO(content_bytes)
            df = pd.read_excel(excel_file)
            text = df.to_string(index=False)
            return text
        except Exception as e:
            print(f"Error extracting XLSX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting XLSX text: {e}")
        return None

def extract_pptx_text_from_bytes(content_bytes):
    try:
        from pptx import Presentation
        pptx_file = io.BytesIO(content_bytes)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except ImportError:
        print("python-pptx not installed. Installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            from pptx import Presentation
            pptx_file = io.BytesIO(content_bytes)
            prs = Presentation(pptx_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PPTX text: {e}")
            return None
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return None

def get_access_token():
    print("Attempting to get access token...")
    
    accounts = app.get_accounts()
    if accounts:
        print("Found existing account")
        result = app.acquire_token_silent(scopes, account=accounts[0])
        if result:
            return result['access_token']
    
    print("Opening browser for login...")
    flow = app.initiate_device_flow(scopes=scopes)
    
    if "user_code" not in flow:
        print("Failed to create device code")
        return None
    
    print(f"Device code: {flow['user_code']}")
    print(f"Login URL: {flow['verification_uri']}")
    
    webbrowser.open(flow['verification_uri'])
    
    result = app.acquire_token_by_device_flow(flow)
    
    if "access_token" in result:
        print("Login successful")
        return result['access_token']
    else:
        print("Login failed")
        print("Error:", result.get('error'))
        print("Description:", result.get('error_description'))
        return None

def get_email_by_id(email_id, headers):
    email_url = f'https://graph.microsoft.com/v1.0/me/messages/{email_id}'
    response = requests.get(email_url, headers=headers)
    
    if response.status_code == 200:
        return response.json()
    else:
        print(f"Error getting email: {response.status_code}")
        print(f"Error: {response.text}")
        return None

def get_attachments_for_email(email_id, headers):
    att_url = f'https://graph.microsoft.com/v1.0/me/messages/{email_id}/attachments'
    response = requests.get(att_url, headers=headers)
    
    if response.status_code == 200:
        return response.json().get('value', [])
    else:
        print(f"Error getting attachments: {response.status_code}")
        return []

def print_file_content_from_bytes(attachment):
    filename = attachment['name']
    content_bytes = base64.b64decode(attachment['contentBytes'])
    file_ext = os.path.splitext(filename.lower())[1]
    
    if is_text_file(filename):
        try:
            content = content_bytes.decode('utf-8')
            if len(content.strip()) == 0:
                print(f"\nFile: {filename}")
                print("="*60)
                print("File is empty or contains no readable text.")
                print("="*60)
                return
            
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        except UnicodeDecodeError:
            try:
                content = content_bytes.decode('latin-1')
                if len(content.strip()) == 0:
                    print(f"\nFile: {filename}")
                    print("="*60)
                    print("File is empty or contains no readable text.")
                    print("="*60)
                    return
                
                print(f"\nContent of {filename}:")
                print("="*60)
                print(content)
                print("="*60)
            except:
                print(f"\nFile: {filename}")
                print("="*60)
                print("Cannot read this file as text.")
                print("="*60)
        except Exception as e:
            print(f"\nFile: {filename}")
            print("="*60)
            print(f"Error reading file: {e}")
            print("="*60)
    
    elif file_ext == '.pdf':
        print(f"\nExtracting text from PDF: {filename}")
        content = extract_pdf_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from PDF.")
            print("="*60)
    
    elif file_ext == '.docx':
        print(f"\nExtracting text from DOCX: {filename}")
        content = extract_docx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from DOCX.")
            print("="*60)
    
    elif file_ext in ['.xlsx', '.xls']:
        print(f"\nExtracting text from Excel: {filename}")
        content = extract_xlsx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from Excel file.")
            print("="*60)
    
    elif file_ext in ['.pptx', '.ppt']:
        print(f"\nExtracting text from PowerPoint: {filename}")
        content = extract_pptx_text_from_bytes(content_bytes)
        if content:
            print(f"\nContent of {filename}:")
            print("="*60)
            print(content)
            print("="*60)
        else:
            print(f"\nFile: {filename}")
            print("="*60)
            print("Could not extract text from PowerPoint file.")
            print("="*60)
    
    else:
        print(f"\nFile: {filename}")
        print("="*60)
        print("This file type is not supported for text extraction.")
        print("="*60)

def main():
    if not client_id or not tenant_id:
        print("Error: CLIENT_ID and TENANT_ID must be set in .env file")
        print("Please create a .env file with the following content:")
        print("CLIENT_ID=your_client_id_here")
        print("TENANT_ID=your_tenant_id_here")
        return
    
    access_token = get_access_token()
    
    if not access_token:
        print("Failed to get access token")
        return
    
    headers = {
        'Authorization': f'Bearer {access_token}',
        'Content-Type': 'application/json'
    }
    
    print("\n" + "="*60)
    print("Email Attachment Viewer")
    print("="*60)
    
    email_id = input("Enter the email ID: ").strip()
    
    if not email_id:
        print("Email ID is required!")
        return
    
    print(f"\nFetching email with ID: {email_id}")
    
    email_data = get_email_by_id(email_id, headers)
    
    if not email_data:
        print("Email not found!")
        return
    
    print(f"\nEmail Details:")
    print(f"Subject: {email_data.get('subject', 'No subject')}")
    print(f"From: {email_data.get('from', {}).get('emailAddress', {}).get('address', 'Unknown')}")
    print(f"Date: {email_data.get('receivedDateTime', 'Unknown')}")
    
    attachments = get_attachments_for_email(email_id, headers)
    
    if not attachments:
        print("\nNo attachments found in this email.")
        return
    
    print(f"\nFound {len(attachments)} attachments:")
    
    allowed_attachments = []
    other_attachments = []
    
    for i, att in enumerate(attachments):
        if att['@odata.type'] == '#microsoft.graph.fileAttachment':
            if is_allowed_file(att['name']):
                file_type = "Text" if is_text_file(att['name']) else "Document"
                allowed_attachments.append(att)
                print(f"{len(allowed_attachments)}. {att['name']} ({file_type} file)")
            else:
                other_attachments.append(att)
                print(f"   {att['name']} (Skipped - not a document)")
    
    if not allowed_attachments:
        print("\nNo document attachments found. Only document files are supported.")
        return
    
    print(f"\nDocument attachments available: {len(allowed_attachments)}")
    print("Note: The script will extract text directly from attachments without saving files.")
    print("Required libraries will be installed automatically if needed.")
    
    while True:
        try:
            choice = input(f"\nSelect a file to view (1-{len(allowed_attachments)}) or 'q' to quit: ").strip()
            
            if choice.lower() == 'q':
                print("Goodbye!")
                break
            
            file_index = int(choice) - 1
            
            if 0 <= file_index < len(allowed_attachments):
                selected_attachment = allowed_attachments[file_index]
                print_file_content_from_bytes(selected_attachment)
                
                view_another = input("\nView another file? (y/n): ").strip().lower()
                if view_another != 'y':
                    break
            else:
                print(f"Please enter a number between 1 and {len(allowed_attachments)}")
        except ValueError:
            print("Please enter a valid number")
        except KeyboardInterrupt:
            print("\nGoodbye!")
            break

if __name__ == "__main__":
    main() 